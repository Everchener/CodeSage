"""
文档加载模块 - 支持多种文件格式的文本提取

支持格式:
- .pdf  - PDF文档 (pdfminer.six)
- .txt  - 纯文本
- .md   - Markdown
- .docx - Word文档 (python-docx)
- .xlsx, .xls - Excel (pandas)
- .pptx - PowerPoint (python-pptx)
- .html - HTML文档 (beautifulsoup4)
- .csv  - CSV文件 (pandas)
"""

import os
import logging
from io import StringIO
from typing import Optional

logger = logging.getLogger(__name__)


def extract_text(filepath: str) -> str:
    """支持多种文件格式的文本提取

    Args:
        filepath: 文件路径

    Returns:
        提取的文本内容，失败返回空字符串
    """
    file_ext = os.path.splitext(filepath)[1].lower()

    if file_ext == '.pdf':
        return _extract_pdf(filepath)
    elif file_ext in ['.txt', '.md']:
        return _extract_text(filepath)
    elif file_ext == '.docx':
        return _extract_docx(filepath)
    elif file_ext in ['.xlsx', '.xls']:
        return _extract_excel(filepath)
    elif file_ext == '.pptx':
        return _extract_pptx(filepath)
    elif file_ext == '.html':
        return _extract_html(filepath)
    elif file_ext == '.csv':
        return _extract_csv(filepath)
    else:
        logger.warning(f"不支持的文件格式: {file_ext}")
        return ""


def _extract_pdf(filepath: str) -> str:
    """提取PDF文本"""
    try:
        from pdfminer.high_level import extract_text_to_fp
        output = StringIO()
        with open(filepath, 'rb') as file:
            extract_text_to_fp(file, output)
        return output.getvalue()
    except Exception as e:
        logger.error(f"PDF提取失败: {e}")
        return ""


def _extract_text(filepath: str) -> str:
    """提取纯文本或Markdown"""
    try:
        with open(filepath, 'r', encoding='utf-8') as file:
            return file.read()
    except UnicodeDecodeError:
        # 尝试其他编码
        try:
            with open(filepath, 'r', encoding='gbk') as file:
                return file.read()
        except Exception as e:
            logger.error(f"文本文件读取失败: {e}")
            return ""
    except Exception as e:
        logger.error(f"文本文件读取失败: {e}")
        return ""


def _extract_docx(filepath: str) -> str:
    """提取Word文档文本"""
    try:
        from docx import Document
        doc = Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])
    except ImportError:
        logger.error("处理Word文档需要安装python-docx库")
        return ""
    except Exception as e:
        logger.error(f"Word文档提取失败: {e}")
        return ""


def _extract_excel(filepath: str) -> str:
    """提取Excel文档文本"""
    try:
        import pandas as pd
        text = ""
        xl = pd.ExcelFile(filepath)
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name)
            text += f"工作表: {sheet_name}\n"
            text += df.to_string(index=False) + "\n\n"
        return text
    except ImportError:
        logger.error("处理Excel文件需要安装pandas和openpyxl库")
        return ""
    except Exception as e:
        logger.error(f"Excel文档提取失败: {e}")
        return ""


def _extract_pptx(filepath: str) -> str:
    """提取PowerPoint文本"""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except ImportError:
        logger.error("处理PPT文件需要安装python-pptx库")
        return ""
    except Exception as e:
        logger.error(f"PPT文档提取失败: {e}")
        return ""


def _extract_html(filepath: str) -> str:
    """提取HTML文档文本"""
    try:
        from bs4 import BeautifulSoup
        with open(filepath, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file.read(), 'lxml')
        # 移除script和style标签
        for script in soup(["script", "style"]):
            script.decompose()
        text = soup.get_text(separator='\n')
        # 清理多余空白
        lines = (line.strip() for line in text.splitlines())
        text = '\n'.join(line for line in lines if line)
        return text
    except ImportError:
        logger.error("处理HTML文件需要安装beautifulsoup4和lxml库")
        return ""
    except Exception as e:
        logger.error(f"HTML文档提取失败: {e}")
        return ""


def _extract_csv(filepath: str) -> str:
    """提取CSV文件文本"""
    try:
        import pandas as pd
        df = pd.read_csv(filepath)
        return df.to_string(index=False)
    except ImportError:
        logger.error("处理CSV文件需要安装pandas库")
        return ""
    except Exception as e:
        logger.error(f"CSV文件提取失败: {e}")
        return ""


def is_supported(filepath: str) -> bool:
    """检查文件格式是否支持"""
    ext = os.path.splitext(filepath)[1].lower()
    supported = ['.pdf', '.txt', '.md', '.docx', '.xlsx', '.xls', '.pptx', '.html', '.csv']
    return ext in supported


# 下面开始是 UTF-8 版本的覆盖实现。
# 这些定义会替换上面的旧实现，同时保持对外模块路径不变。
from dataclasses import dataclass
from typing import Callable

from .config import SUPPORTED_EXTENSIONS as CONFIGURED_SUPPORTED_EXTENSIONS
from .doc_extractors import (
    extract_docx_text,
    extract_excel_text,
    extract_html_text,
    extract_markdown_text,
    extract_pdf_text,
    extract_plain_text,
    extract_pptx_text,
)
from .doc_extractors.base import normalize_text


Extractor = Callable[[str], str]


@dataclass(frozen=True)
class LoaderSpec:
    """描述文档提取器及其旧版兜底实现。"""

    name: str
    extractor: Extractor
    fallback: Extractor | None = None


def _extract_legacy_pdf(filepath: str) -> str:
    try:
        from pdfminer.high_level import extract_text_to_fp

        output = StringIO()
        with open(filepath, "rb") as file:
            extract_text_to_fp(file, output)
        return output.getvalue()
    except Exception:
        logger.exception("旧版 PDF 提取失败：%s", filepath)
        return ""


def _extract_legacy_text(filepath: str) -> str:
    try:
        with open(filepath, "r", encoding="utf-8") as file:
            return file.read()
    except UnicodeDecodeError:
        try:
            with open(filepath, "r", encoding="gb18030") as file:
                return file.read()
        except Exception:
            logger.exception("旧版文本提取失败：%s", filepath)
            return ""
    except Exception:
        logger.exception("旧版文本提取失败：%s", filepath)
        return ""


def _extract_legacy_docx(filepath: str) -> str:
    try:
        from docx import Document

        document = Document(filepath)
        return "\n".join(paragraph.text for paragraph in document.paragraphs)
    except Exception:
        logger.exception("旧版 DOCX 提取失败：%s", filepath)
        return ""


def _extract_legacy_excel(filepath: str) -> str:
    try:
        import pandas as pd

        workbook = pd.ExcelFile(filepath)
        text_blocks = []
        for sheet_name in workbook.sheet_names:
            frame = workbook.parse(sheet_name)
            text_blocks.append(f"工作表：{sheet_name}\n{frame.to_string(index=False)}")
        return "\n\n".join(text_blocks)
    except Exception:
        logger.exception("旧版表格文档提取失败：%s", filepath)
        return ""


def _extract_legacy_csv(filepath: str) -> str:
    try:
        import pandas as pd

        frame = pd.read_csv(filepath)
        return frame.to_string(index=False)
    except Exception:
        logger.exception("旧版 CSV 提取失败：%s", filepath)
        return ""


def _extract_legacy_pptx(filepath: str) -> str:
    try:
        from pptx import Presentation

        presentation = Presentation(filepath)
        texts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        logger.exception("旧版 PPTX 提取失败：%s", filepath)
        return ""


def _extract_legacy_html(filepath: str) -> str:
    try:
        from bs4 import BeautifulSoup

        with open(filepath, "r", encoding="utf-8") as file:
            soup = BeautifulSoup(file.read(), "lxml")
        for tag in soup(["script", "style"]):
            tag.decompose()
        lines = (line.strip() for line in soup.get_text(separator="\n").splitlines())
        return "\n".join(line for line in lines if line)
    except Exception:
        logger.exception("旧版 HTML 提取失败：%s", filepath)
        return ""


LOADERS: dict[str, LoaderSpec] = {
    ".pdf": LoaderSpec("doc_extractors.pdf", extract_pdf_text, _extract_legacy_pdf),
    ".txt": LoaderSpec("doc_extractors.text", extract_plain_text, _extract_legacy_text),
    ".md": LoaderSpec("doc_extractors.markdown", extract_markdown_text, _extract_legacy_text),
    ".docx": LoaderSpec("doc_extractors.docx", extract_docx_text, _extract_legacy_docx),
    ".xlsx": LoaderSpec("doc_extractors.excel", extract_excel_text, _extract_legacy_excel),
    ".xls": LoaderSpec("doc_extractors.excel", extract_excel_text, _extract_legacy_excel),
    ".pptx": LoaderSpec("doc_extractors.pptx", extract_pptx_text, _extract_legacy_pptx),
    ".html": LoaderSpec("doc_extractors.html", extract_html_text, _extract_legacy_html),
    ".csv": LoaderSpec("doc_extractors.excel", extract_excel_text, _extract_legacy_csv),
}


def _get_extension(filepath: str) -> str:
    return os.path.splitext(filepath)[1].lower()


def _get_loader(filepath: str) -> LoaderSpec | None:
    return LOADERS.get(_get_extension(filepath))


def resolve_loader_name(filepath: str) -> str:
    """返回指定文件路径对应的提取器名称。"""
    loader = _get_loader(filepath)
    return loader.name if loader else "unsupported"


def _run_extractor(extractor: Extractor, filepath: str, *, loader_name: str) -> str:
    extracted_text = extractor(filepath)
    normalized = normalize_text(extracted_text or "")
    if normalized:
        logger.debug("文档已通过 %s 提取：%s", loader_name, filepath)
    return normalized


def extract_text(filepath: str) -> str:
    """从受支持的文档路径中提取文本。"""
    loader = _get_loader(filepath)
    if loader is None:
        logger.warning("不支持的文档格式：%s", filepath)
        return ""

    try:
        text = _run_extractor(loader.extractor, filepath, loader_name=loader.name)
        if text:
            return text
        if loader.fallback is not None:
            logger.warning(
                "%s 对 %s 返回空文本，回退到旧版提取器",
                loader.name,
                filepath,
            )
            return _run_extractor(loader.fallback, filepath, loader_name=f"{loader.name}.legacy")
    except Exception:
        logger.exception("主提取器 %s 执行失败：%s", loader.name, filepath)
        if loader.fallback is not None:
            return _run_extractor(loader.fallback, filepath, loader_name=f"{loader.name}.legacy")
    return ""


def is_supported(filepath: str) -> bool:
    """检查文件扩展名是否受支持。"""
    return _get_extension(filepath) in set(CONFIGURED_SUPPORTED_EXTENSIONS)
