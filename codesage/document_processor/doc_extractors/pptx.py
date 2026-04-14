"""按稳定形状顺序提取 PPTX 内容的辅助函数。"""

from __future__ import annotations

from .base import join_sections, prefix_heading, table_rows_to_text


def _sort_shapes(shapes) -> list:
    return sorted(
        shapes,
        key=lambda shape: (
            (shape.top if shape.top is not None else 0) // 10,
            shape.left if shape.left is not None else 0,
        ),
    )


def _paragraph_text(paragraph) -> str:
    text = paragraph.text.strip()
    if not text:
        return ""
    if paragraph.level > 0:
        return f"{'  ' * paragraph.level}- {text}"
    return text


def _table_text(table) -> str:
    rows = []
    for row in table.rows:
        rows.append([cell.text for cell in row.cells])
    return table_rows_to_text(rows)


def _extract_shape_text(shape) -> str:
    if getattr(shape, "has_text_frame", False):
        lines = [_paragraph_text(paragraph) for paragraph in shape.text_frame.paragraphs]
        return "\n".join(line for line in lines if line)

    if getattr(shape, "has_table", False):
        return _table_text(shape.table)

    if hasattr(shape, "shapes"):
        nested = [_extract_shape_text(child) for child in _sort_shapes(shape.shapes)]
        return "\n".join(text for text in nested if text)

    if hasattr(shape, "text"):
        return str(shape.text).strip()

    return ""


def extract_text(filepath: str) -> str:
    """从 PPTX 文件中提取按顺序排列的幻灯片内容。"""
    from pptx import Presentation

    presentation = Presentation(filepath)
    slides: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        slide_lines = [_extract_shape_text(shape) for shape in _sort_shapes(slide.shapes)]
        slide_text = join_sections(line for line in slide_lines if line)
        if slide_text:
            slides.append(join_sections([prefix_heading(f"第 {index} 页幻灯片", 1), slide_text]))

    return join_sections(slides)
